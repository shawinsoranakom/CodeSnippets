def convert_file(self, file_path: Path, **kwargs) -> ConvertResult:
        try:
            from pptx import Presentation
            from pptx.shapes.picture import Picture
        except ImportError:
            raise RuntimeError(
                "PPTX conversion requires python-pptx: pip install paddleocr[doc2md]"
            )
        try:
            from pptx.enum.shapes import MSO_SHAPE_TYPE
        except ImportError:
            MSO_SHAPE_TYPE = None  # type: ignore[assignment]

        prs = Presentation(str(file_path))
        slides_md = []
        images: dict = {}
        image_counter = [0]
        slide_width = prs.slide_width
        self._Picture = Picture
        self._MSO_SHAPE_TYPE = MSO_SHAPE_TYPE

        for slide in prs.slides:
            slide_parts = []

            # Process all shapes
            for shape in slide.shapes:
                self._process_shape(
                    shape, slide_parts, images, image_counter, slide_width, slide.part
                )

            # Handle math formulas inside mc:AlternateContent elements
            # (python-pptx doesn't expose these as Shape objects)
            for alt_content in slide._element.iter(_MC_ALT):
                # Only look at mc:Choice (the preferred rendering path)
                choice = alt_content.find(_MC_CHOICE)
                if choice is None:
                    continue
                for para_elem in choice.iter(_A_P):
                    if _paragraph_has_math(para_elem):
                        math_items = _extract_math_from_paragraph(para_elem)
                        for latex in math_items:
                            slide_parts.append(f"$$\n{latex}\n$$")

            # Speaker notes
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    slide_parts.append(f"\n> **Notes**: {notes_text}")

            # Group parts by content type and separate groups with blank lines
            # to prevent HTML blocks from consuming adjacent list items
            groups: list[list[str]] = []
            for part in slide_parts:
                kind = _classify_part(part)
                if groups and _classify_part(groups[-1][0]) == kind:
                    groups[-1].append(part)
                else:
                    groups.append([part])

            slides_md.append("\n\n".join("\n".join(g) for g in groups))

        md_text = "\n\n---\n\n".join(slides_md)

        return ConvertResult(
            markdown=md_text,
            images=images,
            metadata={
                "format": "PPTX",
                "slide_count": len(prs.slides),
            },
        )